import json
import boto3
import os
import re
import logging
from urllib.parse import unquote_plus
import tempfile
from io import BytesIO
import time
from botocore.exceptions import ClientError

# Configure logging first
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Document processing libraries with better error handling
try:
    from docx import Document
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError as e:
    logger.warning(f"DOCX import failed: {str(e)}")
    Document = None
    DOCX_AVAILABLE = False

try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError as e:
    logger.warning(f"pypdf import failed: {str(e)}")
    pypdf = None
    PYPDF_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError as e:
    logger.warning(f"openpyxl import failed: {str(e)}")
    load_workbook = None
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError as e:
    logger.warning(f"python-pptx import failed: {str(e)}")
    Presentation = None
    PPTX_AVAILABLE = False

# Initialize AWS clients
s3 = boto3.client('s3')

# Configuration constants
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB limit
MAX_CONFIG_SIZE = 1 * 1024 * 1024  # 1MB limit for config
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'doc', 'xlsx', 'xls', 'csv', 'pptx', 'ppt'}
MAX_REPLACEMENTS = 100  # Maximum number of replacement rules
BATCH_SIZE = 5  # Maximum files to process in one batch
BATCH_TIMEOUT = 45  # Maximum seconds for batch processing

# Retry configuration
MAX_RETRIES = 3
BASE_BACKOFF = 1  # seconds
MAX_BACKOFF = 30  # seconds

# Environment variables
INPUT_BUCKET = os.environ['INPUT_BUCKET']
OUTPUT_BUCKET = os.environ['OUTPUT_BUCKET']
QUARANTINE_BUCKET = os.environ['QUARANTINE_BUCKET']
CONFIG_BUCKET = os.environ['CONFIG_BUCKET']

# Windows compatibility mode (for ChatGPT compatibility)
WINDOWS_MODE = os.environ.get('WINDOWS_MODE', 'true').lower() == 'true'

# Global cache for config
_config_cache = None
_config_last_modified = None

# Default fallback patterns if config not available
DEFAULT_REPLACEMENTS = [
    {"find": "REPLACE_CLIENT_NAME", "replace": "Company X"},
    {"find": "ACME Corporation", "replace": "[REDACTED]"},
    {"find": "TechnoSoft", "replace": "[REDACTED]"}
]

# Common PII regex patterns
PII_PATTERNS = {
    "ssn": {
        "pattern": r"\b(?:\d{3}-\d{2}-\d{4}|\d{3}\s\d{2}\s\d{4}|\d{9})\b",
        "replace": "[SSN]",
        "description": "Social Security Number"
    },
    "credit_card": {
        "pattern": r"\b(?:\d{4}[\s\-]?){3}\d{4}\b",
        "replace": "[CREDIT_CARD]",
        "description": "Credit Card Number"
    },
    "phone": {
        "pattern": r"\b(?:\+?1[\s\-\.]?)?\(?\d{3}\)?[\s\-\.]?\d{3}[\s\-\.]?\d{4}\b",
        "replace": "[PHONE]",
        "description": "Phone Number (US/Canada)"
    },
    "email": {
        "pattern": r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b",
        "replace": "[EMAIL]",
        "description": "Email Address"
    },
    "ip_address": {
        "pattern": r"\b(?:(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)\.){3}(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)\b",
        "replace": "[IP_ADDRESS]",
        "description": "IPv4 Address"
    },
    "drivers_license": {
        "pattern": r"\b[A-Z]{1,2}\d{5,8}\b",
        "replace": "[DRIVERS_LICENSE]",
        "description": "Driver's License (various states)"
    }
}

def get_user_info_from_key(key):
    """Extract user info from S3 key if it follows user prefix pattern"""
    # Check if key follows pattern: users/{user_id}/...
    if key.startswith('users/'):
        parts = key.split('/', 2)
        if len(parts) >= 3:
            return {
                'user_id': parts[1],
                'file_path': parts[2]
            }
    return None

def validate_file(bucket, key):
    """
    Validate file before processing
    """
    try:
        # Check file extension
        file_path = key
        user_info = get_user_info_from_key(key)
        if user_info:
            file_path = user_info['file_path']
            
        file_ext = file_path.lower().split('.')[-1]
        if file_ext not in ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {file_ext}. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}")
        
        # Check file size
        response = s3.head_object(Bucket=bucket, Key=key)
        file_size = response['ContentLength']
        
        if file_size > MAX_FILE_SIZE:
            raise ValueError(f"File too large: {file_size} bytes. Maximum allowed: {MAX_FILE_SIZE} bytes")
        
        if file_size == 0:
            raise ValueError("File is empty")
        
        return True
    except s3.exceptions.NoSuchKey:
        raise ValueError(f"File not found: {key}")
    except Exception as e:
        logger.error(f"File validation error: {str(e)}")
        raise

def get_default_config():
    """Return a clean default configuration for new users"""
    return {
        "replacements": [],
        "case_sensitive": False,
        "patterns": {
            "ssn": False,
            "credit_card": False,
            "phone": False,
            "email": False,
            "ip_address": False,
            "drivers_license": False
        }
    }

def get_redaction_config(user_id=None):
    """Load user-specific configuration from S3 bucket with caching"""
    global _config_cache, _config_last_modified
    
    # Always use user-specific config key
    if not user_id:
        logger.warning("No user_id provided, using default config")
        return get_default_config()
    
    config_key = f'configs/users/{user_id}/config.json'
    
    try:
        # Try user-specific config first
        if user_id:
            try:
                response = s3.head_object(Bucket=CONFIG_BUCKET, Key=config_key)
                last_modified = response['LastModified']
                
                # Use cached config if available and not modified
                cache_key = f"user_{user_id}"
                if isinstance(_config_cache, dict) and cache_key in _config_cache:
                    if _config_last_modified and _config_last_modified.get(cache_key) == last_modified:
                        logger.info(f"Using cached configuration for user {user_id}")
                        return _config_cache[cache_key]
                
                # Load fresh user config
                response = s3.get_object(Bucket=CONFIG_BUCKET, Key=config_key)
                config_content = response['Body'].read()
                
                # Validate size
                if len(config_content) > MAX_CONFIG_SIZE:
                    raise ValueError(f"Configuration file too large: {len(config_content)} bytes")
                
                config = json.loads(config_content)
                
                # Update cache
                if not isinstance(_config_cache, dict):
                    _config_cache = {}
                if not isinstance(_config_last_modified, dict):
                    _config_last_modified = {}
                
                _config_cache[cache_key] = config
                _config_last_modified[cache_key] = last_modified
                
                logger.info(f"Loaded user-specific configuration for {user_id} with {len(config.get('replacements', []))} rules")
                return config
                
            except s3.exceptions.NoSuchKey:
                logger.info(f"No user-specific config for {user_id}, returning default config")
                return get_default_config()
        
        # This should not be reached given the changes above
        logger.warning("Unexpected: No configuration found, using defaults")
        return get_default_config()
        
    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON in configuration: {str(e)}")
        return get_default_config()
    except Exception as e:
        logger.error(f"Error loading configuration: {str(e)}")
        return get_default_config()

def lambda_handler(event, context):
    """
    Main Lambda handler for document processing with batch support and user isolation
    """
    results = []
    start_time = time.time()
    
    try:
        # Get files to process
        files_to_process = []
        for record in event['Records']:
            bucket = record['s3']['bucket']['name']
            key = unquote_plus(record['s3']['object']['key'])
            files_to_process.append((bucket, key))
        
        # Log batch start
        logger.info(json.dumps({
            'event': 'BATCH_START',
            'batch_size': len(files_to_process),
            'request_id': context.aws_request_id
        }))
        
        # Process files in batches
        processed_count = 0
        for i in range(0, len(files_to_process), BATCH_SIZE):
            batch = files_to_process[i:i + BATCH_SIZE]
            
            # Check timeout
            elapsed = time.time() - start_time
            if elapsed > BATCH_TIMEOUT:
                logger.warning(json.dumps({
                    'event': 'BATCH_TIMEOUT',
                    'processed': processed_count,
                    'remaining': len(files_to_process) - processed_count,
                    'elapsed': elapsed
                }))
                break
            
            # Process batch (config will be loaded per user)
            batch_results = process_file_batch(batch, None, context)
            results.extend(batch_results)
            processed_count += len(batch_results)
        
        # Log batch completion
        logger.info(json.dumps({
            'event': 'BATCH_COMPLETE',
            'processed': processed_count,
            'total': len(files_to_process),
            'duration': time.time() - start_time,
            'request_id': context.aws_request_id
        }))
        
        return {
            'statusCode': 200,
            'body': json.dumps({
                'message': f'Processed {processed_count} files',
                'results': results
            })
        }
        
    except Exception as e:
        logger.error(json.dumps({
            'event': 'BATCH_ERROR',
            'error': str(e),
            'type': type(e).__name__,
            'request_id': context.aws_request_id
        }))
        
        # If any files were processed, return partial success
        if results:
            return {
                'statusCode': 207,  # Multi-status
                'body': json.dumps({
                    'message': f'Partial batch processing: {len(results)} succeeded',
                    'error': str(e),
                    'results': results
                })
            }
        
        return {
            'statusCode': 500,
            'body': json.dumps({
                'error': str(e),
                'type': type(e).__name__
            })
        }

def process_file_batch(batch, config, context):
    """Process a batch of files with error isolation and user-specific configs"""
    results = []
    
    for bucket, key in batch:
        try:
            # Extract user info from key
            user_info = get_user_info_from_key(key)
            
            # Load user-specific config for this file
            if user_info:
                user_config = get_redaction_config(user_info['user_id'])
            else:
                # Fall back to global config for non-user files
                user_config = get_redaction_config()
            
            # Validate config
            validate_config(user_config)
            
            # Process individual file with user-specific config
            result = process_single_file(bucket, key, user_config)
            results.append({
                'key': key,
                'status': 'success',
                'result': result,
                'user_id': user_info['user_id'] if user_info else 'global'
            })
            
        except Exception as e:
            error_msg = str(e)
            logger.error(json.dumps({
                'event': 'FILE_ERROR',
                'file': key,
                'error': error_msg,
                'type': type(e).__name__,
                'request_id': context.aws_request_id
            }))
            
            # Try to quarantine the file
            try:
                quarantine_document(bucket, key, error_msg)
                results.append({
                    'key': key,
                    'status': 'quarantined',
                    'error': error_msg
                })
            except Exception as q_error:
                results.append({
                    'key': key,
                    'status': 'failed',
                    'error': error_msg,
                    'quarantine_error': str(q_error)
                })
    
    return results

def process_single_file(bucket, key, config):
    """Process a single file with user isolation support"""
    logger.info(f"Processing file: s3://{bucket}/{key}")
    
    # Extract user info if present
    user_info = get_user_info_from_key(key)
    
    # Validate file
    validate_file(bucket, key)
    
    # Get file extension
    file_path = user_info['file_path'] if user_info else key
    file_ext = file_path.lower().split('.')[-1]
    
    # Process based on file type
    if file_ext == 'txt':
        process_text_file(bucket, key, config, user_info)
    elif file_ext == 'pdf':
        process_pdf_file(bucket, key, config, user_info)
    elif file_ext in ['docx', 'doc']:
        process_docx_file(bucket, key, config, user_info)
    elif file_ext in ['xlsx', 'xls']:
        process_xlsx_file(bucket, key, config, user_info)
    elif file_ext == 'csv':
        process_csv_file(bucket, key, config, user_info)
    elif file_ext in ['pptx', 'ppt']:
        process_pptx_file(bucket, key, config, user_info)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")
    
    # Delete original file from input bucket after successful processing
    delete_processed_file(bucket, key)
    
    return {'processed': True, 'file_type': file_ext}

def strip_urls_preserve_text(text):
    """
    Strip URLs from text while preserving the link text.
    Handles various formats:
    - HTML: <a href="url">text</a> -> text
    - Markdown: [text](url) -> text
    - Plain URLs: http://example.com -> (removed)
    - Email links: <a href="mailto:email">text</a> -> text
    """
    import re
    
    # Strip HTML anchor tags but keep the text
    # Matches <a href="...">text</a> and replaces with just text
    text = re.sub(r'<a\s+(?:[^>]*?\s+)?href=["\'](?:[^"\']*)["\'][^>]*>(.*?)</a>', r'\1', text, flags=re.IGNORECASE | re.DOTALL)
    
    # Strip Markdown links but keep the text
    # Matches [text](url) and replaces with just text
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    
    # Strip standalone URLs (http, https, ftp, etc.)
    # This regex matches URLs that are not part of HTML or Markdown syntax
    url_pattern = r'(?<!["\'\(])\b(?:https?|ftp|ftps)://[^\s<>"{}|\\^`\[\]]+(?!["\'\)])'
    text = re.sub(url_pattern, '', text)
    
    # Strip www URLs without protocol
    www_pattern = r'(?<!["\'\(/@])\bwww\.[^\s<>"{}|\\^`\[\]]+(?!["\'\)])'
    text = re.sub(www_pattern, '', text)
    
    # Clean up any double spaces left after URL removal
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()

def apply_redaction_rules(text, config):
    """Apply redaction rules to text content including pattern-based PII detection"""
    if not config:
        return text, False
    
    redacted = False
    
    # Strip URLs first while preserving link text
    processed_text = strip_urls_preserve_text(text)
    if processed_text != text:
        redacted = True
    
    case_sensitive = config.get('case_sensitive', False)
    
    # Apply conditional rules first (content-based redaction)
    conditional_rules = config.get('conditional_rules', [])
    for rule in conditional_rules:
        # Skip disabled rules
        if not rule.get('enabled', True):
            continue
            
        rule_name = rule.get('name', 'Unnamed Rule')
        trigger = rule.get('trigger', {})
        trigger_contains = trigger.get('contains', [])
        trigger_case_sensitive = trigger.get('case_sensitive', False)
        
        # Check if any trigger words exist in the text
        trigger_matched = False
        for trigger_word in trigger_contains:
            if not trigger_word:
                continue
                
            if trigger_case_sensitive:
                if trigger_word in text:
                    trigger_matched = True
                    break
            else:
                if trigger_word.lower() in text.lower():
                    trigger_matched = True
                    break
        
        # If trigger matched, apply the rule's replacements
        if trigger_matched:
            logger.info(f"Conditional rule '{rule_name}' triggered")
            rule_replacements = rule.get('replacements', [])
            for replacement in rule_replacements:
                find_text = replacement.get('find', '')
                replace_text = replacement.get('replace', '[REDACTED]')
                
                if not find_text:
                    continue
                
                # Use rule's case sensitivity setting or inherit from trigger
                rule_case_sensitive = replacement.get('case_sensitive', trigger_case_sensitive)
                if rule_case_sensitive:
                    pattern = re.escape(find_text)
                else:
                    pattern = re.escape(find_text)
                    pattern = f"(?i){pattern}"
                
                # Apply replacement
                if re.search(pattern, processed_text):
                    redacted = True
                    processed_text = re.sub(pattern, replace_text, processed_text)
                    logger.info(f"Applied conditional redaction: '{find_text}' -> '{replace_text}'")
    
    # Apply global text-based replacements
    replacements = config.get('replacements', [])
    for replacement in replacements:
        find_text = replacement.get('find', '')
        replace_text = replacement.get('replace', '[REDACTED]')
        
        if not find_text:
            continue
        
        # Create pattern based on case sensitivity
        if case_sensitive:
            pattern = re.escape(find_text)
        else:
            pattern = re.escape(find_text)
            pattern = f"(?i){pattern}"
        
        # Check if pattern exists in text
        if re.search(pattern, processed_text):
            redacted = True
            processed_text = re.sub(pattern, replace_text, processed_text)
            logger.info(f"Applied redaction: '{find_text}' -> '{replace_text}'")
    
    # Apply pattern-based PII detection if enabled
    patterns = config.get('patterns', {})
    if patterns:
        for pattern_name, enabled in patterns.items():
            if enabled and pattern_name in PII_PATTERNS:
                pii_config = PII_PATTERNS[pattern_name]
                pattern = pii_config['pattern']
                replace = pii_config['replace']
                
                # Check if pattern exists in text
                if re.search(pattern, processed_text):
                    redacted = True
                    processed_text = re.sub(pattern, replace, processed_text)
                    logger.info(f"Applied PII pattern '{pattern_name}': {pii_config['description']}")
    
    # Normalize text output for better compatibility
    processed_text = normalize_text_output(processed_text)
    
    return processed_text, redacted

def apply_redaction_with_count(text, config):
    """Apply redaction rules and return both the redacted text and count of replacements"""
    if not config:
        return text, 0
    
    # Strip URLs first while preserving link text
    processed_text = strip_urls_preserve_text(text)
    replacement_count = 0
    if processed_text != text:
        replacement_count += 1  # Count URL stripping as one replacement
    
    case_sensitive = config.get('case_sensitive', False)
    
    # Apply conditional rules first (content-based redaction)
    conditional_rules = config.get('conditional_rules', [])
    for rule in conditional_rules:
        # Skip disabled rules
        if not rule.get('enabled', True):
            continue
            
        rule_name = rule.get('name', 'Unnamed Rule')
        trigger = rule.get('trigger', {})
        trigger_contains = trigger.get('contains', [])
        trigger_case_sensitive = trigger.get('case_sensitive', False)
        
        # Check if any trigger words exist in the text
        trigger_matched = False
        for trigger_word in trigger_contains:
            if not trigger_word:
                continue
                
            if trigger_case_sensitive:
                if trigger_word in text:
                    trigger_matched = True
                    break
            else:
                if trigger_word.lower() in text.lower():
                    trigger_matched = True
                    break
        
        # If trigger matched, apply the rule's replacements
        if trigger_matched:
            logger.info(f"Conditional rule '{rule_name}' triggered")
            rule_replacements = rule.get('replacements', [])
            for replacement in rule_replacements:
                find_text = replacement.get('find', '')
                replace_text = replacement.get('replace', '[REDACTED]')
                
                if not find_text:
                    continue
                
                # Use rule's case sensitivity setting or inherit from trigger
                rule_case_sensitive = replacement.get('case_sensitive', trigger_case_sensitive)
                if rule_case_sensitive:
                    pattern = re.escape(find_text)
                else:
                    pattern = re.escape(find_text)
                    pattern = f"(?i){pattern}"
                
                # Count and apply replacements
                matches = list(re.finditer(pattern, processed_text))
                if matches:
                    replacement_count += len(matches)
                    processed_text = re.sub(pattern, replace_text, processed_text)
                    logger.info(f"Applied conditional redaction: '{find_text}' -> '{replace_text}' ({len(matches)} times)")
    
    # Apply global text-based replacements
    replacements = config.get('replacements', [])
    for replacement in replacements:
        find_text = replacement.get('find', '')
        replace_text = replacement.get('replace', '[REDACTED]')
        
        if not find_text:
            continue
        
        # Create pattern based on case sensitivity
        if case_sensitive:
            pattern = re.escape(find_text)
        else:
            pattern = re.escape(find_text)
            pattern = f"(?i){pattern}"
        
        # Count and apply replacements
        matches = list(re.finditer(pattern, processed_text))
        if matches:
            replacement_count += len(matches)
            processed_text = re.sub(pattern, replace_text, processed_text)
            logger.info(f"Applied redaction: '{find_text}' -> '{replace_text}' ({len(matches)} times)")
    
    # Apply pattern-based PII detection if enabled
    patterns = config.get('patterns', {})
    if patterns:
        for pattern_name, enabled in patterns.items():
            if enabled and pattern_name in PII_PATTERNS:
                pii_config = PII_PATTERNS[pattern_name]
                pattern = pii_config['pattern']
                replace = pii_config['replace']
                
                # Count and apply replacements
                matches = list(re.finditer(pattern, processed_text))
                if matches:
                    replacement_count += len(matches)
                    processed_text = re.sub(pattern, replace, processed_text)
                    logger.info(f"Applied PII pattern '{pattern_name}': {pii_config['description']} ({len(matches)} times)")
    
    # Normalize text output for better compatibility
    processed_text = normalize_text_output(processed_text)
    
    return processed_text, replacement_count

def normalize_text_output(text, windows_mode=None):
    """Normalize text for Markdown format compatibility with ChatGPT
    
    This function:
    1. Converts line endings based on mode (Unix or Windows)
    2. Replaces special UTF-8 characters with ASCII equivalents
    3. Ensures consistent encoding for plain text output
    4. Outputs clean text that works with ChatGPT's file upload
    
    Note: Files are saved as .md extension but contain plain text (no markdown formatting)
    This is a workaround for ChatGPT's .txt/.log file upload issues
    """
    if not text:
        return text
    
    # Use global setting if not specified
    if windows_mode is None:
        windows_mode = WINDOWS_MODE
    
    # Log that normalization is running
    logger.info("Running text normalization - input length: %d, windows_mode: %s", len(text), windows_mode)
    
    # First normalize to Unix line endings
    text = text.replace('\r\n', '\n')
    text = text.replace('\r', '\n')
    
    # Remove BOM if present
    if text.startswith('\ufeff'):
        text = text[1:]
        logger.info("Removed BOM from text")
    
    # Replace common special characters that cause issues
    replacements = {
        # Curly quotes to straight quotes
        '\u2018': "'",  # Left single quotation mark
        '\u2019': "'",  # Right single quotation mark
        '\u201C': '"',  # Left double quotation mark
        '\u201D': '"',  # Right double quotation mark
        
        # Dashes and hyphens
        '\u2013': '-',  # En dash
        '\u2014': '--', # Em dash
        '\u2015': '--', # Horizontal bar
        
        # Other common problematic characters
        '\u2026': '...', # Horizontal ellipsis
        '\u00A0': ' ',   # Non-breaking space
        '\u2022': '*',   # Bullet point
        '\u2023': '>',   # Triangular bullet
        '\u25CF': '*',   # Black circle
        '\u25CB': 'o',   # White circle
        '\u2192': '->',  # Rightwards arrow
        '\u2190': '<-',  # Leftwards arrow
        '\u2194': '<->', # Left right arrow
        
        # Mathematical symbols
        '\u00D7': 'x',   # Multiplication sign
        '\u00F7': '/',   # Division sign
        '\u00B1': '+/-', # Plus-minus sign
        '\u2264': '<=',  # Less than or equal to
        '\u2265': '>=',  # Greater than or equal to
        '\u2260': '!=',  # Not equal to
        
        # Other quotation marks
        '\u00AB': '<<',  # Left-pointing double angle quotation mark
        '\u00BB': '>>',  # Right-pointing double angle quotation mark
        '\u201A': ',',   # Single low-9 quotation mark
        '\u201E': ',,',  # Double low-9 quotation mark
        
        # Additional spaces and special characters
        '\u202F': ' ',   # Narrow no-break space
        '\u2009': ' ',   # Thin space
        '\u200A': ' ',   # Hair space
        '\u2008': ' ',   # Punctuation space
        '\u205F': ' ',   # Medium mathematical space
        '\u3000': ' ',   # Ideographic space
        '\u00AD': '',    # Soft hyphen
        '\u2011': '-',   # Non-breaking hyphen
        '\u2212': '-',   # Minus sign
        '\uFEFF': '',    # Zero-width no-break space (BOM)
    }
    
    for old_char, new_char in replacements.items():
        text = text.replace(old_char, new_char)
    
    # Convert tabs to spaces for better compatibility
    text = text.replace('\t', '    ')
    
    # Remove any remaining non-printable characters (except newlines)
    # This preserves ASCII 32-126, plus newline (10)
    cleaned_text = ''.join(
        char if (32 <= ord(char) <= 126) or char == '\n' 
        else ' ' if ord(char) > 126 else ''
        for char in text
    )
    
    # Final pass: ensure absolutely no non-ASCII characters remain
    # Use 'ignore' instead of 'replace' to avoid question marks
    final_text = cleaned_text.encode('ascii', 'ignore').decode('ascii')
    
    # Clean up multiple spaces
    import re
    final_text = re.sub(r' +', ' ', final_text)
    final_text = re.sub(r'\n +', '\n', final_text)
    final_text = re.sub(r' +\n', '\n', final_text)
    
    # Convert to Windows line endings if requested
    if windows_mode:
        final_text = final_text.replace('\n', '\r\n')
        # Optionally add UTF-8 BOM for Windows compatibility
        # Note: BOM is not ASCII, so only add if specifically needed
        # final_text = '\ufeff' + final_text
    
    # Log completion
    logger.info("Text normalization complete - output length: %d, is ASCII: %s, line_ending: %s", 
                len(final_text), 
                all(ord(c) < 128 for c in final_text),
                "CRLF" if windows_mode else "LF")
    
    return final_text

def apply_filename_redaction(filename, config):
    """Apply redaction rules to file names, preserving the file extension"""
    try:
        # Extract base name and extension
        if '.' in filename:
            base_name, extension = filename.rsplit('.', 1)
        else:
            base_name = filename
            extension = 'md'  # Default to .md if no extension
        
        # Apply redaction to base name only
        processed_name, redacted = apply_redaction_rules(base_name, config)
        
        # Preserve the original extension (already set by file type processors)
        processed_filename = f"{processed_name}.{extension}"
            
        if redacted:
            logger.info(f"Applied filename redaction: {filename} -> {processed_filename}")
            
        return processed_filename, redacted
        
    except Exception as e:
        logger.error(f"Error applying filename redaction: {str(e)}")
        return filename, False

def validate_config(config):
    """Validate configuration structure and limits"""
    if not isinstance(config, dict):
        raise ValueError("Configuration must be a dictionary")
    
    # Validate replacements
    replacements = config.get('replacements', [])
    if not isinstance(replacements, list):
        raise ValueError("Replacements must be a list")
    
    if len(replacements) > MAX_REPLACEMENTS:
        raise ValueError(f"Too many replacement rules: {len(replacements)} (max: {MAX_REPLACEMENTS})")
    
    for i, replacement in enumerate(replacements):
        if not isinstance(replacement, dict):
            raise ValueError(f"Replacement {i} must be a dictionary")
        if 'find' not in replacement:
            raise ValueError(f"Replacement {i} missing 'find' field")
        if not replacement['find']:
            raise ValueError(f"Replacement {i} has empty 'find' field")
    
    # Validate patterns if present
    patterns = config.get('patterns', {})
    if patterns and not isinstance(patterns, dict):
        raise ValueError("Patterns must be a dictionary")
    
    # Validate pattern names
    for pattern_name in patterns:
        if pattern_name not in PII_PATTERNS:
            logger.warning(f"Unknown pattern name: {pattern_name}")

def exponential_backoff_retry(func):
    """Execute function with exponential backoff retry"""
    last_exception = None
    
    for attempt in range(MAX_RETRIES):
        try:
            return func()
        except ClientError as e:
            if e.response['Error']['Code'] in ['RequestTimeout', 'ServiceUnavailable', 'ThrottlingException']:
                last_exception = e
                backoff = min(BASE_BACKOFF * (2 ** attempt), MAX_BACKOFF)
                logger.warning(f"Attempt {attempt + 1} failed, retrying in {backoff}s: {str(e)}")
                time.sleep(backoff)
            else:
                raise
        except Exception as e:
            raise
    
    raise last_exception

def upload_processed_document(key, content, metadata=None, config=None, user_info=None):
    """Upload processed document to output bucket with retry logic and user isolation"""
    # Extract filename from key
    if user_info:
        # Remove user prefix from key for processing
        filename = user_info['file_path']
    else:
        filename = key
    
    # Apply filename redaction if config provided
    if config:
        redacted_filename, filename_redacted = apply_filename_redaction(filename, config)
        if metadata is None:
            metadata = {}
        metadata['filename_redacted'] = str(filename_redacted)
    else:
        redacted_filename = filename
    
    # Construct the output key with user prefix if needed
    if user_info:
        processed_key = f"processed/users/{user_info['user_id']}/{redacted_filename}"
    else:
        processed_key = f"processed/{redacted_filename}"
    
    base_metadata = {
        'processing-status': 'completed',
        'original-key': key
    }
    
    if user_info:
        base_metadata['user-id'] = user_info['user_id']
    
    if metadata:
        base_metadata.update(metadata)
    
    # Handle both bytes and file-like objects
    if hasattr(content, 'read'):
        body = content.read()
    else:
        body = content
    
    def _upload():
        s3.put_object(
            Bucket=OUTPUT_BUCKET,
            Key=processed_key,
            Body=body,
            ServerSideEncryption='AES256',
            Metadata=base_metadata
        )
    
    try:
        exponential_backoff_retry(_upload)
        logger.info(f"Uploaded processed document: s3://{OUTPUT_BUCKET}/{processed_key}")
        return True
    except Exception as e:
        logger.error(f"Error uploading processed document after {MAX_RETRIES} attempts: {str(e)}")
        raise

def process_text_file(bucket, key, config, user_info=None):
    """Process text files with comprehensive error handling"""
    try:
        # Download file with retry
        def _download():
            return s3.get_object(Bucket=bucket, Key=key)
        
        response = exponential_backoff_retry(_download)
        content = response['Body'].read().decode('utf-8', errors='replace')
        
        # Apply redaction rules
        processed_content, redacted = apply_redaction_rules(content, config)
        
        # Upload processed document
        upload_processed_document(key, processed_content.encode('utf-8'), 
                                {'redacted': str(redacted)}, config, user_info)
        
        return True
        
    except Exception as e:
        logger.error(f"Error processing text file {key}: {str(e)}")
        raise

def process_pdf_file(bucket, key, config, user_info=None):
    """Process PDF files by extracting text and redacting"""
    if not PYPDF_AVAILABLE:
        raise ImportError("pypdf library not available for PDF processing")
    
    try:
        # Download file
        def _download():
            return s3.get_object(Bucket=bucket, Key=key)
        
        response = exponential_backoff_retry(_download)
        pdf_content = response['Body'].read()
        
        # Extract text from PDF
        pdf_file = BytesIO(pdf_content)
        reader = pypdf.PdfReader(pdf_file)
        
        all_text = []
        for page_num, page in enumerate(reader.pages):
            try:
                text = page.extract_text()
                if text.strip():
                    all_text.append(f"--- Page {page_num + 1} ---\n{text}")
            except Exception as e:
                logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                all_text.append(f"--- Page {page_num + 1} ---\n[Error extracting text]")
        
        # Combine all text
        full_text = '\n\n'.join(all_text)
        
        if not full_text.strip():
            raise ValueError("No text content found in PDF")
        
        # Apply redaction rules
        processed_text, redacted = apply_redaction_rules(full_text, config)
        
        # Save as markdown file (change extension to .md for ChatGPT compatibility)
        file_path = user_info['file_path'] if user_info else key
        text_key = file_path.rsplit('.', 1)[0] + '.md'
        
        # Update user_info with the new filename for proper handling
        if user_info:
            updated_user_info = user_info.copy()
            updated_user_info['file_path'] = text_key
            text_key = f"users/{user_info['user_id']}/{text_key}"
        else:
            updated_user_info = None
        
        # Upload processed document as text
        upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                {'redacted': str(redacted), 'converted_from': 'pdf'}, config, updated_user_info)
        
        return True
        
    except Exception as e:
        logger.error(f"Error processing PDF file {key}: {str(e)}")
        raise

def process_docx_file(bucket, key, config, user_info=None):
    """Process DOCX files with fallback methods"""
    try:
        # Download file
        def _download():
            return s3.get_object(Bucket=bucket, Key=key)
        
        response = exponential_backoff_retry(_download)
        docx_content = response['Body'].read()
        
        # Try ZIP-based extraction first (doesn't require lxml)
        text_content = extract_docx_text_zip(docx_content)
        
        if not text_content and DOCX_AVAILABLE:
            # Fallback to python-docx if available
            text_content = extract_docx_text_library(docx_content)
        
        if not text_content:
            raise ValueError("No text content found in DOCX")
        
        # Apply redaction rules
        processed_text, redacted = apply_redaction_rules(text_content, config)
        
        # Save as markdown file
        file_path = user_info['file_path'] if user_info else key
        text_key = file_path.rsplit('.', 1)[0] + '.md'
        
        # Update user_info with the new filename for proper handling
        if user_info:
            updated_user_info = user_info.copy()
            updated_user_info['file_path'] = text_key
            text_key = f"users/{user_info['user_id']}/{text_key}"
        else:
            updated_user_info = None
        
        # Upload processed document
        upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                {'redacted': str(redacted), 'converted_from': 'docx'}, config, updated_user_info)
        
        return True
        
    except Exception as e:
        logger.error(f"Error processing DOCX file {key}: {str(e)}")
        raise

def extract_docx_text_zip(docx_content):
    """Extract text from DOCX using ZIP method (no dependencies)"""
    import zipfile
    import xml.etree.ElementTree as ET
    
    try:
        with zipfile.ZipFile(BytesIO(docx_content)) as zip_file:
            # Extract document.xml
            with zip_file.open('word/document.xml') as doc_xml:
                tree = ET.parse(doc_xml)
                root = tree.getroot()
                
                # Extract all text elements
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                }
                
                paragraphs = []
                for paragraph in root.findall('.//w:p', namespaces):
                    texts = []
                    for text_elem in paragraph.findall('.//w:t', namespaces):
                        if text_elem.text:
                            texts.append(text_elem.text)
                    if texts:
                        paragraphs.append(''.join(texts))
                
                return '\n'.join(paragraphs)
                
    except Exception as e:
        logger.warning(f"ZIP extraction failed: {str(e)}")
        return None

def extract_docx_text_library(docx_content):
    """Extract text from DOCX using python-docx library"""
    if not DOCX_AVAILABLE:
        return None
        
    try:
        doc = Document(BytesIO(docx_content))
        paragraphs = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append('\t'.join(row_text))
        
        return '\n'.join(paragraphs)
        
    except Exception as e:
        logger.warning(f"python-docx extraction failed: {str(e)}")
        return None

def process_csv_file(bucket, key, config, user_info=None):
    """Process CSV files by reading and applying redaction rules"""
    try:
        # Download file
        response = s3.get_object(Bucket=bucket, Key=key)
        csv_content = response['Body'].read()
        
        # Try to decode with UTF-8, fallback to latin-1
        try:
            text = csv_content.decode('utf-8')
        except UnicodeDecodeError:
            text = csv_content.decode('latin-1')
        
        # Apply redaction rules to the CSV content
        processed_text, redacted = apply_redaction_rules(text, config)
        
        # Normalize text for Windows compatibility
        processed_text = normalize_text_output(processed_text)
        
        # Change file extension to .md for ChatGPT compatibility
        file_path = user_info['file_path'] if user_info else key
        text_key = file_path.rsplit('.', 1)[0] + '.md'
        
        # Update user_info with the new filename for proper handling
        if user_info:
            updated_user_info = user_info.copy()
            updated_user_info['file_path'] = text_key
            text_key = f"users/{user_info['user_id']}/{text_key}"
        else:
            updated_user_info = None
        
        # Upload processed document
        upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                {'redacted': str(redacted), 'converted_from': 'csv'}, config, updated_user_info)
        
        return True
        
    except Exception as e:
        logger.error(f"Error processing CSV file {key}: {str(e)}")
        raise

def process_xlsx_file(bucket, key, config, user_info=None):
    """Process XLSX files by converting first sheet to CSV format"""
    if not OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl library not available for XLSX processing")
    
    try:
        # Download file
        def _download():
            return s3.get_object(Bucket=bucket, Key=key)
        
        response = exponential_backoff_retry(_download)
        xlsx_content = response['Body'].read()
        
        # Load workbook
        workbook = load_workbook(BytesIO(xlsx_content), read_only=True, data_only=True)
        
        # Get sheet count and names
        sheet_count = len(workbook.worksheets)
        sheet_names = [sheet.title for sheet in workbook.worksheets]
        
        # Process only the first worksheet
        if sheet_count == 0:
            raise ValueError("Workbook has no sheets")
            
        first_sheet = workbook.worksheets[0]
        csv_rows = []
        
        # Add header comment about sheet count if multiple sheets
        if sheet_count > 1:
            csv_rows.append(f"# Workbook contains {sheet_count} sheets. Showing sheet 1 of {sheet_count}: '{first_sheet.title}'")
            csv_rows.append(f"# Other sheets: {', '.join(sheet_names[1:])}")
            csv_rows.append("")  # Empty line after comments
        
        # Convert first sheet to CSV format
        for row in first_sheet.iter_rows():
            row_values = []
            for cell in row:
                value = cell.value if cell.value is not None else ""
                # Escape values that contain commas, quotes, or newlines
                value_str = str(value)
                if ',' in value_str or '"' in value_str or '\n' in value_str:
                    # Escape quotes by doubling them
                    value_str = value_str.replace('"', '""')
                    # Wrap in quotes
                    value_str = f'"{value_str}"'
                row_values.append(value_str)
            if row_values:
                csv_rows.append(','.join(row_values))
        
        # Combine all rows
        full_text = '\n'.join(csv_rows)
        
        # Apply redaction rules
        processed_text, redacted = apply_redaction_rules(full_text, config)
        
        # Save as CSV file
        file_path = user_info['file_path'] if user_info else key
        text_key = file_path.rsplit('.', 1)[0] + '.csv'
        
        # Update user_info with the new filename for proper handling
        if user_info:
            updated_user_info = user_info.copy()
            updated_user_info['file_path'] = text_key
            text_key = f"users/{user_info['user_id']}/{text_key}"
        else:
            updated_user_info = None
        
        # Upload processed document as CSV
        metadata = {
            'redacted': str(redacted), 
            'converted_from': 'xlsx',
            'sheet_count': str(sheet_count),
            'sheets_included': '1',
            'first_sheet_name': first_sheet.title
        }
        if sheet_count > 1:
            metadata['omitted_sheets'] = ', '.join(sheet_names[1:])
            
        upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                metadata, config, updated_user_info)
        
        return True
        
    except Exception as e:
        logger.error(f"Error processing XLSX file {key}: {str(e)}")
        raise

def process_pptx_file(bucket, key, config, user_info=None):
    """Process PowerPoint files by extracting text from all slides"""
    # Use simple extraction if python-pptx is not available
    if not PPTX_AVAILABLE:
        logger.info("Using simple PPTX extraction method")
        from pptx_handler import process_pptx_simple
        
        try:
            # Download file
            def _download():
                return s3.get_object(Bucket=bucket, Key=key)
            
            response = exponential_backoff_retry(_download)
            pptx_content = response['Body'].read()
            
            # Process with simple handler
            processed_text, redacted = process_pptx_simple(pptx_content, config)
            
            # Change file extension to .md for ChatGPT compatibility
            file_path = user_info['file_path'] if user_info else key
            text_key = file_path.rsplit('.', 1)[0] + '.md'
            
            # Update user_info with the new filename for proper handling
            if user_info:
                updated_user_info = user_info.copy()
                updated_user_info['file_path'] = text_key
                text_key = f"users/{user_info['user_id']}/{text_key}"
            else:
                updated_user_info = None
            
            # Upload processed document
            metadata = {
                'redacted': str(redacted),
                'converted_from': 'pptx',
                'extraction_method': 'simple'
            }
            
            upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                    metadata, config, updated_user_info)
            
            logger.info(f"Successfully processed PPTX file using simple method: {key}")
            return
        except Exception as e:
            logger.error(f"Failed to process PPTX with simple method: {str(e)}")
            raise
    
    try:
        # Download file
        def _download():
            return s3.get_object(Bucket=bucket, Key=key)
        
        response = exponential_backoff_retry(_download)
        pptx_content = response['Body'].read()
        
        # Create a Presentation object
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=True) as tmp_file:
            tmp_file.write(pptx_content)
            tmp_file.seek(0)
            
            presentation = Presentation(tmp_file.name)
            
            # Extract text from all slides
            extracted_text = []
            slide_count = len(presentation.slides)
            
            for idx, slide in enumerate(presentation.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {idx} ---")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                    
                    # Check for tables
                    if shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    extracted_text.extend(slide_text)
                    extracted_text.append("")  # Empty line between slides
            
            # Add metadata header
            header = f"# PowerPoint Document ({slide_count} slides)\n\n"
            full_text = header + "\n".join(extracted_text)
            
            # Apply redaction rules
            processed_text, redacted = apply_redaction_rules(full_text, config)
            
            # Normalize text for Windows compatibility
            processed_text = normalize_text_for_windows(processed_text)
            
            # Change file extension to .md for ChatGPT compatibility
            file_path = user_info['file_path'] if user_info else key
            text_key = file_path.rsplit('.', 1)[0] + '.md'
            
            # Update user_info with the new filename for proper handling
            if user_info:
                updated_user_info = user_info.copy()
                updated_user_info['file_path'] = text_key
                text_key = f"users/{user_info['user_id']}/{text_key}"
            else:
                updated_user_info = None
            
            # Upload processed document
            metadata = {
                'redacted': str(redacted),
                'converted_from': 'pptx',
                'slide_count': str(slide_count)
            }
            
            upload_processed_document(text_key, processed_text.encode('utf-8'), 
                                    metadata, config, updated_user_info)
            
            return True
        
    except Exception as e:
        logger.error(f"Error processing PPTX file {key}: {str(e)}")
        raise

def quarantine_document(bucket, key, reason):
    """Move document to quarantine bucket with retry logic and user isolation"""
    copy_source = {'Bucket': bucket, 'Key': key}
    
    # Check for user prefix
    user_info = get_user_info_from_key(key)
    if user_info:
        quarantine_key = f"quarantine/users/{user_info['user_id']}/{user_info['file_path']}"
    else:
        quarantine_key = f"quarantine/{key}"
    
    def _quarantine():
        s3.copy_object(
            CopySource=copy_source,
            Bucket=QUARANTINE_BUCKET,
            Key=quarantine_key,
            ServerSideEncryption='AES256',
            Metadata={
                'quarantine-reason': reason[:255],  # S3 metadata value limit
                'original-bucket': bucket,
                'original-key': key,
                'user-id': user_info['user_id'] if user_info else 'none'
            }
        )
    
    try:
        exponential_backoff_retry(_quarantine)
        logger.info(f"Document quarantined: s3://{QUARANTINE_BUCKET}/{quarantine_key}, reason: {reason}")
    except Exception as e:
        logger.error(f"Error quarantining document after {MAX_RETRIES} attempts: {str(e)}")
        raise

def delete_processed_file(bucket, key):
    """Delete file from input bucket after successful processing"""
    try:
        def _delete():
            s3.delete_object(Bucket=bucket, Key=key)
        
        exponential_backoff_retry(_delete)
        logger.info(f"Deleted processed file: s3://{bucket}/{key}")
    except Exception as e:
        logger.error(f"Error deleting processed file: {str(e)}")
        # Don't raise exception as file was processed successfully
        
def create_success_marker(bucket, key):
    """Create a success marker file (optional for tracking)"""
    try:
        marker_key = f"_success/{key}.done"
        s3.put_object(
            Bucket=OUTPUT_BUCKET,
            Key=marker_key,
            Body=json.dumps({
                'processed_at': time.time(),
                'original_bucket': bucket,
                'original_key': key
            }),
            ServerSideEncryption='AES256'
        )
    except Exception as e:
        logger.warning(f"Failed to create success marker: {str(e)}")
        # Don't fail the processing for marker creation